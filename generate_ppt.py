"""
ChanStock README PPT 生成器
Usage: python generate_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import sys, os

# ── 颜色常量 ──────────────────────────────────────────────
DARK_BG     = RGBColor(0x06, 0x08, 0x0C)   # #06080c 深色背景
GREEN       = RGBColor(0x22, 0xC5, 0x5E)   # #22c55e 主绿色
TEAL        = RGBColor(0x2D, 0xD4, 0xBF)   # #2dd4bf 渐变青
CYAN        = RGBColor(0x38, 0xBD, 0xF8)   # #38bdf8 蓝
ORANGE      = RGBColor(0xF5, 0x9E, 0x0B)   # #f59e0b 橙色
RED         = RGBColor(0xEF, 0x44, 0x44)   # #ef4444 红色
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0x94, 0xA3, 0xB8)   # #94a3b8 浅灰文字
MID_GRAY    = RGBColor(0x47, 0x48, 0x56)   # #474856 中灰
CARD_BG     = RGBColor(0x0E, 0x16, 0x21)   # #0e1621 卡片背景
SLIDE_BG    = RGBColor(0x07, 0x09, 0x0F)   # #07090f

prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)

BLANK_LAYOUT = prs.slide_layouts[6]  # 空白布局

# ── 辅助函数 ──────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK_LAYOUT)

def fill_bg(slide, color=DARK_BG):
    """填充全幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # 调整圆角
    shape.adjustments[0] = 0.12
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, font_name="微软雅黑"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_multiline_text(slide, lines, left, top, width, height,
                       font_size=14, color=LIGHT_GRAY, line_spacing=1.5,
                       font_name="微软雅黑"):
    """lines: list of (text, bold, color) tuples"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, col = item, False, color
        else:
            text, bold, col = item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = col
        run.font.name = font_name
    return txBox

def add_tag(slide, text, left, top, bg_color, text_color=CYAN, font_size=11):
    """小标签（药丸形状）"""
    shape = slide.shapes.add_shape(
        5,  # rounded rect
        Inches(left), Inches(top), Inches(len(text) * 0.1 + 0.25), Inches(0.32)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.9

    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = "Consolas"

def add_divider(slide, left, top, width, color=GREEN, height=0.02):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ════════════════════════════════════════════════════════════
#  第 1 页：封面（Hero Slide）
# ════════════════════════════════════════════════════════════
s1 = add_slide()
fill_bg(s1, SLIDE_BG)

# 顶部绿色渐变条
add_rect(s1, 0, 0, 16, 0.06, GREEN)

# 中心品牌区背景
brand_bg = s1.shapes.add_shape(
    1, Inches(0), Inches(2.5), Inches(16), Inches(4)
)
brand_bg.fill.solid()
brand_bg.fill.fore_color.rgb = RGBColor(0x07, 0x09, 0x0F)
brand_bg.line.fill.background()

# 主标题
add_text(s1, "ChanStock", 0, 2.7, 16, 1.2,
         font_size=72, bold=True,
         color=GREEN, align=PP_ALIGN.CENTER)

# 副标题
add_text(s1, "缠论智能股票分析系统", 0, 4.0, 16, 0.6,
         font_size=22, bold=False,
         color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 分隔线
add_divider(s1, 5.5, 4.7, 5, TEAL, 0.025)

# 4 个功能亮点
features = [
    ("📐", "缠论结构识别", "分型 · 笔 · 线段 · 中枢"),
    ("🎯", "买卖点判定", "一买/二买/三买 · 一卖/二卖/三卖"),
    ("🤖", "AI 背驰判断", "MACD 面积比较 · 多级别共振"),
    ("📊", "多级别 K 线", "1分~日线 · 周线 · 月线"),
]
for i, (icon, title, desc) in enumerate(features):
    x = 0.7 + i * 3.8
    card = add_rounded_rect(s1, x, 5.0, 3.4, 1.1, CARD_BG)
    add_text(s1, icon, x + 0.1, 5.08, 0.5, 0.5, font_size=20)
    add_text(s1, title, x + 0.6, 5.08, 2.6, 0.35, font_size=12, bold=True, color=WHITE)
    add_text(s1, desc,  x + 0.6, 5.42, 2.6, 0.35, font_size=10, color=LIGHT_GRAY)

# 底部技术栈标签
tags = ["Python 3.10+", "FastAPI", "AKShare", "Vue 3", "TypeScript", "ECharts"]
total_w = sum(len(t) * 0.095 + 0.3 for t in tags)
start_x = (16 - total_w) / 2
cur_x = start_x
for t in tags:
    w = len(t) * 0.095 + 0.3
    bg = CARD_BG
    col = CYAN if t in ("FastAPI", "Vue 3", "TypeScript") else \
          ORANGE if t in ("AKShare", "Python 3.10+") else GREEN
    shape = s1.shapes.add_shape(
        5, Inches(cur_x), Inches(6.3), Inches(w), Inches(0.32)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = col
    shape.line.width = Pt(0.5)
    shape.adjustments[0] = 0.9
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = t
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = col
    run.font.name = "Consolas"
    cur_x += w + 0.12

# 版本号
add_text(s1, "v0.1.0", 14.5, 8.5, 1.2, 0.35, font_size=10,
         color=MID_GRAY, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════
#  第 2 页：项目概览
# ════════════════════════════════════════════════════════════
s2 = add_slide()
fill_bg(s2, SLIDE_BG)
add_rect(s2, 0, 0, 16, 0.06, GREEN)

add_text(s2, "项目概览", 0.5, 0.25, 15, 0.7,
         font_size=32, bold=True, color=WHITE)
add_divider(s2, 0.5, 0.95, 3, GREEN, 0.03)

add_text(s2,
    "ChanStock 是一款面向 A 股的智能技术分析工具，"
    "核心逻辑基于缠中说禅理论，通过程序化识别分型、笔、线段、"
    "中枢等结构，判断趋势、背驰与买卖点，并结合 AI 大模型（DeepSeek / Gemini）"
    "输出可操作的投资建议。",
    0.5, 1.2, 15, 1.2, font_size=15, color=LIGHT_GRAY)

# 三列目标用户
users = [
    ("🧑‍🎓", "缠论学习者", "希望将缠论应用于实战，\n借助系统快速识别笔、\n线段、中枢", GREEN),
    ("📈", "技术分析爱好者", "希望综合 MA、MACD、\nRSI 等指标与\n缠论结构", ORANGE),
    ("💻", "量化研究开发者", "基于缠论数据接口\n做二次策略开发", CYAN),
]
for i, (icon, title, desc, col) in enumerate(users):
    x = 0.6 + i * 5.1
    add_rounded_rect(s2, x, 2.5, 4.5, 2.8, CARD_BG)
    add_text(s2, icon, x + 0.2, 2.65, 1, 0.7, font_size=32)
    add_text(s2, title, x + 0.2, 3.35, 4, 0.45, font_size=15, bold=True, color=WHITE)
    # 分隔线
    add_divider(s2, x + 0.2, 3.8, 4, col, 0.02)
    add_text(s2, desc, x + 0.2, 3.9, 4, 1.2, font_size=12, color=LIGHT_GRAY)

# 免责
add_divider(s2, 0.5, 5.6, 15, MID_GRAY)
add_text(s2,
    "⚠️  免责声明：本系统仅供技术研究与学习使用，不构成任何投资建议。股票投资有风险，入市需谨慎。",
    0.5, 5.7, 15, 0.4, font_size=11, color=MID_GRAY)

# ════════════════════════════════════════════════════════════
#  第 3 页：功能特性总览
# ════════════════════════════════════════════════════════════
s3 = add_slide()
fill_bg(s3, SLIDE_BG)
add_rect(s3, 0, 0, 16, 0.06, GREEN)

add_text(s3, "功能特性", 0.5, 0.25, 15, 0.7,
         font_size=32, bold=True, color=WHITE)
add_divider(s3, 0.5, 0.95, 3, GREEN, 0.03)

# 6 个功能卡片 2行3列
features_grid = [
    ("📐", "缠论结构识别", [
        "自动识别顶/底分型",
        "笔检测（≥5根K线）",
        "线段识别（≥3笔）",
        "中枢识别（三段重叠）",
        "走势：涨/跌/盘/震",
    ], GREEN),
    ("🎯", "买卖点判定", [
        "一买：下跌背驰点",
        "二买：一买后回调",
        "三买：突破中枢回踩",
        "一卖/二卖/三卖",
        "置信度量化评估",
    ], RED),
    ("📊", "K 线可视化", [
        "MA5/20/60 均线",
        "笔（红涨绿跌）",
        "线段（黄/橙粗线）",
        "中枢（紫色矩形）",
        "副图：MACD/RSI/SKDJ",
    ], ORANGE),
    ("🤖", "AI 策略增强", [
        "MACD 面积背驰检测",
        "多级别共振分析",
        "规则策略引擎",
        "LLM 自然语言分析",
        "止损止盈建议",
    ], CYAN),
    ("🔍", "股票数据", [
        "代码/名称模糊搜索",
        "实时行情（实时）",
        "多周期K线",
        "热门股票/板块",
        "自选股管理",
    ], RGBColor(0xA8, 0x55, 0xF7)),
    ("📐", "支撑阻力位", [
        "中枢上下沿（强）",
        "线段高低点（次强）",
        "笔高低点（近期）",
        "买卖点对应价格",
        "历史K线高低价",
    ], RGBColor(0xF4, 0x72, 0x2E)),
]

for i, (icon, title, items, col) in enumerate(features_grid):
    row, col_idx = i // 3, i % 3
    x = 0.5 + col_idx * 5.1
    y = 1.15 + row * 3.7
    add_rounded_rect(s3, x, y, 4.9, 3.4, CARD_BG)
    # 左侧色条
    add_rect(s3, x, y + 0.08, 0.06, 3.24, col)
    add_text(s3, icon + "  " + title, x + 0.2, y + 0.1, 4.4, 0.5,
             font_size=14, bold=True, color=WHITE)
    add_divider(s3, x + 0.2, y + 0.55, 4.4, col, 0.015)
    for j, item in enumerate(items):
        add_text(s3, "• " + item, x + 0.2, y + 0.65 + j * 0.5, 4.4, 0.45,
                  font_size=11, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════════
#  第 4 页：缠论核心算法
# ════════════════════════════════════════════════════════════
s4 = add_slide()
fill_bg(s4, SLIDE_BG)
add_rect(s4, 0, 0, 16, 0.06, TEAL)

add_text(s4, "缠论核心算法", 0.5, 0.25, 15, 0.7,
         font_size=32, bold=True, color=WHITE)
add_divider(s4, 0.5, 0.95, 4, TEAL, 0.03)

# 左侧：算法步骤流程
steps = [
    (GREEN, "分型检测", "顶分型（∧）/ 底分型（∨）\n包含关系处理：取高高 / 取低低"),
    (TEAL,  "笔识别",   "顶分型 + 底分型 = 一笔\n默认 ≥5 根 K 线过滤噪音"),
    (CYAN,  "线段识别", "≥3 笔构成，代表次级别走势\n需被反向线段破坏"),
    (ORANGE,"中枢识别", "三段以上重叠区域\n多空博弈均衡区间"),
    (RED,   "背驰判断", "底背驰：价格新低但 MACD 未新低\n顶背驰：价格新高但 MACD 未新高"),
]
for i, (col, title, desc) in enumerate(steps):
    y = 1.2 + i * 1.52
    add_rounded_rect(s4, 0.5, y, 8.5, 1.35, CARD_BG)
    add_rect(s4, 0.5, y + 0.05, 0.06, 1.25, col)
    add_text(s4, title, 0.72, y + 0.1, 2, 0.45, font_size=14, bold=True, color=col)
    add_text(s4, desc,  0.72, y + 0.52, 7.8, 0.7, font_size=11, color=LIGHT_GRAY)
    if i < len(steps) - 1:
        add_text(s4, "▼", 4.4, y + 1.25, 0.6, 0.3,
                 font_size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)

# 右侧：买卖点详解
add_rounded_rect(s4, 9.5, 1.2, 6, 4.3, CARD_BG)
add_text(s4, "买卖点体系", 9.7, 1.3, 5.5, 0.5,
         font_size=15, bold=True, color=WHITE)
add_divider(s4, 9.7, 1.78, 5.5, GREEN, 0.02)

buys = [
    ("一买", "下跌趋势背驰点\n力度 < 前段×0.8 且价格创新低", GREEN),
    ("二买", "一买后回调低点\n不低于一买点（×0.95容差）", TEAL),
    ("三买", "突破中枢后回踩\n低点不跌入中枢上沿", CYAN),
]
for i, (name, desc, col) in enumerate(buys):
    y = 1.9 + i * 0.75
    add_rounded_rect(s4, 9.7, y, 1.1, 0.55, col)
    t = add_text(s4, name, 9.72, y + 0.08, 1.05, 0.4,
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s4, desc, 10.9, y + 0.04, 4.4, 0.55, font_size=10, color=LIGHT_GRAY)

add_divider(s4, 9.7, 4.15, 5.5, MID_GRAY)
sells = [
    ("一卖", "上涨趋势背驰点\n力度 < 前段×0.8 且价格创新高", RED),
    ("二卖", "一卖后反弹高点\n不超一卖点（×1.02容差）", ORANGE),
    ("三卖", "跌破中枢后反弹\n高点不突破中枢下沿", RGBColor(0xA8, 0x55, 0xF7)),
]
for i, (name, desc, col) in enumerate(sells):
    y = 4.28 + i * 0.38
    add_rounded_rect(s4, 9.7, y, 1.1, 0.28, col)
    t = add_text(s4, name, 9.72, y + 0.04, 1.05, 0.22,
                 font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s4, desc, 10.9, y + 0.04, 4.4, 0.3, font_size=10, color=LIGHT_GRAY)

# 置信度说明
add_text(s4, "置信度加成：背驰加持+15% | 多中枢衰减×0.85 | 趋势加持+10%", 9.7, 5.55, 5.6, 0.35,
         font_size=10, color=MID_GRAY)

# ════════════════════════════════════════════════════════════
#  第 5 页：系统架构
# ════════════════════════════════════════════════════════════
s5 = add_slide()
fill_bg(s5, SLIDE_BG)
add_rect(s5, 0, 0, 16, 0.06, GREEN)

add_text(s5, "系统架构", 0.5, 0.25, 15, 0.7,
         font_size=32, bold=True, color=WHITE)
add_divider(s5, 0.5, 0.95, 3, GREEN, 0.03)

# 前端卡片
add_rounded_rect(s5, 0.5, 1.2, 4.5, 5.5, CARD_BG)
add_rect(s5, 0.5, 1.2, 4.5, 0.55, GREEN)
add_text(s5, "🌐  前端 (Browser)", 0.5, 1.25, 4.5, 0.45,
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_multiline_text(s5, [
    ("Vue 3  + TypeScript", True, WHITE),
    ("ECharts 5.5.x（K线可视化）", False, LIGHT_GRAY),
    ("Pinia（状态管理）", False, LIGHT_GRAY),
    ("Vue Router（路由）", False, LIGHT_GRAY),
    ("Axios（HTTP客户端）", False, LIGHT_GRAY),
    ("", False, LIGHT_GRAY),
    ("视图层", True, CYAN),
    ("• HomeView（首页）", False, LIGHT_GRAY),
    ("• StockView（个股分析）", False, LIGHT_GRAY),
    ("• WatchlistView（自选股）", False, LIGHT_GRAY),
], 0.65, 1.9, 4.1, 4.5, font_size=11)

# 中间箭头
add_text(s5, "⟷", 4.6, 3.5, 1, 0.5, font_size=28, color=MID_GRAY, align=PP_ALIGN.CENTER)

# 后端卡片
add_rounded_rect(s5, 5.7, 1.2, 10, 5.5, CARD_BG)
add_rect(s5, 5.7, 1.2, 10, 0.55, CYAN)
add_text(s5, "⚙️  后端 (Server)", 5.7, 1.25, 10, 0.45,
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

layers = [
    (GREEN, "REST API 层 (FastAPI)", "main.py — /api/stocks/* · /api/chanlun/* · /api/watchlist/*"),
    (TEAL,  "缠论引擎 (ChanlunEngine)", "分型 → 笔 → 线段 → 中枢 → 背驰 → 买卖点"),
    (ORANGE,"策略引擎 (StrategyEngine)", "信号优先级 · 置信度 · 止损止盈计算"),
    (RED,   "AI 层 (LLM Client)", "DeepSeek / Gemini — 自然语言缠论分析"),
    (CYAN,  "数据服务 (AKShare)", "东方财富 / 新浪 / 腾讯 — K线/行情/新闻/板块"),
]
for i, (col, title, desc) in enumerate(layers):
    y = 1.9 + i * 0.96
    add_rounded_rect(s5, 5.85, y, 9.7, 0.8, RGBColor(0x0B, 0x10, 0x1A))
    add_rect(s5, 5.85, y + 0.05, 0.06, 0.7, col)
    add_text(s5, title, 6.05, y + 0.05, 4.5, 0.38,
             font_size=12, bold=True, color=col)
    add_text(s5, desc, 6.05, y + 0.42, 9.2, 0.32,
             font_size=10, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════════
#  第 6 页：安装部署
# ════════════════════════════════════════════════════════════
s6 = add_slide()
fill_bg(s6, SLIDE_BG)
add_rect(s6, 0, 0, 16, 0.06, GREEN)

add_text(s6, "安装部署", 0.5, 0.25, 15, 0.7,
         font_size=32, bold=True, color=WHITE)
add_divider(s6, 0.5, 0.95, 3, GREEN, 0.03)

# 步骤卡片
steps_deploy = [
    ("1", "克隆项目", "bash\ngit clone <repo-url>\ncd stock-chanlun", GREEN),
    ("2", "安装后端", "bash\ncd backend\npython -m venv .venv\n.venv\\Scripts\\activate\npip install -r requirements.txt", TEAL),
    ("3", "启动后端", "bash\n# 推荐（自动处理网络代理）\npython run_server.py\n\n# 或直接 uvicorn\nuvicorn main:app --port 8000", CYAN),
    ("4", "安装前端", "bash\ncd frontend\nnpm install", ORANGE),
    ("5", "启动前端", "bash\nnpm run dev\n# http://localhost:5173", RED),
]

for i, (num, title, code, col) in enumerate(steps_deploy):
    row, c = i // 3, i % 3
    x = 0.5 + c * 5.1
    y = 1.15 + row * 4.0
    add_rounded_rect(s6, x, y, 4.9, 3.6, CARD_BG)
    # 序号圆
    circ = s6.shapes.add_shape(9, Inches(x + 0.15), Inches(y + 0.15),
                                Inches(0.45), Inches(0.45))
    circ.fill.solid(); circ.fill.fore_color.rgb = col
    circ.line.fill.background()
    tf = circ.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = num
    run.font.size = Pt(16); run.font.bold = True
    run.font.color.rgb = WHITE
    add_text(s6, title, x + 0.7, y + 0.2, 3.5, 0.4,
             font_size=14, bold=True, color=WHITE)
    add_rect(s6, x + 0.15, y + 0.62, 4.5, 0.015, col)
    # 代码框
    code_bg = s6.shapes.add_shape(1, Inches(x + 0.15), Inches(y + 0.7),
                                   Inches(4.6), Inches(2.7))
    code_bg.fill.solid(); code_bg.fill.fore_color.rgb = RGBColor(0x0A, 0x0F, 0x17)
    code_bg.line.color.rgb = RGBColor(0x1E, 0x2A, 0x38)
    code_bg.line.width = Pt(0.5)
    add_text(s6, code, x + 0.25, y + 0.75, 4.4, 2.5,
             font_size=9, color=RGBColor(0xA5, 0xB4, 0xC4), font_name="Consolas")

# 环境变量说明
add_rounded_rect(s6, 0.5, 5.3, 15, 1.1, CARD_BG)
add_text(s6, "📋  环境变量（可选）", 0.7, 5.4, 5, 0.4,
         font_size=12, bold=True, color=WHITE)
add_text(s6,
    "在 backend/.env 配置 AI Key，不配置则以纯规则模式运行（所有缠论分析仍然正常工作）：\n"
    "DEEPSEEK_API_KEY=sk-xxxxxxxxxxxxxxxxxxxx    |    GEMINI_API_KEY=AIzaSyxxxxxxxxxxxxxxxxxxxxx",
    0.7, 5.78, 14.5, 0.55, font_size=10, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════════
#  第 7 页：API 接口一览
# ════════════════════════════════════════════════════════════
s7 = add_slide()
fill_bg(s7, SLIDE_BG)
add_rect(s7, 0, 0, 16, 0.06, GREEN)

add_text(s7, "API 接口一览", 0.5, 0.25, 15, 0.7,
         font_size=32, bold=True, color=WHITE)
add_divider(s7, 0.5, 0.95, 4, GREEN, 0.03)
add_text(s7, "基础路径：http://localhost:8000/api", 0.5, 1.05, 15, 0.35,
         font_size=12, color=LIGHT_GRAY)

apis = [
    ("股票数据", [
        ("GET",  "GET /api/stocks/search?q={keyword}",       "股票模糊搜索"),
        ("GET",  "GET /api/stocks/{code}/quote",            "实时行情"),
        ("GET",  "GET /api/stocks/{code}/kline",            "K线数据（level=daily等）"),
        ("GET",  "GET /api/stocks/hot?limit=20",            "热门股票排行"),
        ("GET",  "GET /api/market/overview",                 "市场概览（板块+宽度）"),
        ("GET",  "GET /api/news?limit=8",                   "财经新闻"),
    ]),
    ("缠论分析", [
        ("GET",  "GET /api/chanlun/{code}?level=daily",      "缠论完整分析"),
        ("GET",  "GET /api/chanlun/{code}/ai",              "AI策略信号（可选model=deepseek/gemini）"),
    ]),
    ("自选股管理", [
        ("GET",   "GET    /api/watchlist",                   "获取自选股列表"),
        ("POST",  "POST   /api/watchlist/{code}",            "添加自选股"),
        ("DELETE","DELETE /api/watchlist/{code}",            "删除自选股"),
    ]),
]

cols_x = [0.5, 5.6, 10.7]
for ci, (cat, items) in enumerate(apis):
    x = cols_x[ci]
    add_rounded_rect(s7, x, 1.45, 4.8, 6.4, CARD_BG)
    add_rect(s7, x, 1.45, 4.8, 0.45,
             GREEN if ci == 0 else TEAL if ci == 1 else CYAN)
    add_text(s7, cat, x + 0.15, 1.5, 4.5, 0.38,
             font_size=13, bold=True, color=WHITE)
    for j, (method, path, desc) in enumerate(items):
        y = 2.0 + j * 0.88
        col_m = GREEN if method == "GET" else ORANGE if method == "POST" else RED
        add_rounded_rect(s7, x + 0.1, y, 0.72, 0.3, col_m)
        add_text(s7, method, x + 0.12, y + 0.04, 0.7, 0.24,
                 font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s7, path, x + 0.95, y + 0.03, 3.7, 0.3,
                 font_size=9, color=CYAN, font_name="Consolas")
        add_text(s7, desc, x + 0.1, y + 0.38, 4.5, 0.4,
                 font_size=10, color=LIGHT_GRAY)
        if j < len(items) - 1:
            add_rect(s7, x + 0.1, y + 0.82, 4.5, 0.01, RGBColor(0x1E, 0x2A, 0x38))

# ════════════════════════════════════════════════════════════
#  第 8 页：前端页面布局
# ════════════════════════════════════════════════════════════
s8 = add_slide()
fill_bg(s8, SLIDE_BG)
add_rect(s8, 0, 0, 16, 0.06, GREEN)

add_text(s8, "前端页面说明", 0.5, 0.25, 15, 0.7,
         font_size=32, bold=True, color=WHITE)
add_divider(s8, 0.5, 0.95, 4, GREEN, 0.03)

# 首页
add_rounded_rect(s8, 0.5, 1.15, 4.9, 3.6, CARD_BG)
add_rect(s8, 0.5, 1.15, 4.9, 0.45, GREEN)
add_text(s8, "🏠  首页 /", 0.5, 1.2, 4.9, 0.4,
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_multiline_text(s8, [
    ("• 顶部导航（Logo + 搜索框）", False, LIGHT_GRAY),
    ("• 热门股票排行", False, LIGHT_GRAY),
    ("• 当日热门板块网格", False, LIGHT_GRAY),
    ("• 点击卡片跳转个股分析页", False, LIGHT_GRAY),
    ("", False, LIGHT_GRAY),
    ("路由：/", True, CYAN),
], 0.65, 1.7, 4.5, 2.8, font_size=11)

# 个股分析页
add_rounded_rect(s8, 5.55, 1.15, 4.9, 3.6, CARD_BG)
add_rect(s8, 5.55, 1.15, 4.9, 0.45, CYAN)
add_text(s8, "📈  个股分析页 /stock/{code}", 5.55, 1.2, 4.9, 0.4,
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_multiline_text(s8, [
    ("左栏(240px)：股票代码+行情", False, LIGHT_GRAY),
    ("  级别切换 · 走势判断 · 买卖点列表", False, LIGHT_GRAY),
    ("中栏(自适应)：K线主图+副图", False, LIGHT_GRAY),
    ("  MA均线 · 笔/线段/中枢叠加", False, LIGHT_GRAY),
    ("  MACD · RSI · SKDJ副图", False, LIGHT_GRAY),
    ("右栏(280px)：AI策略卡片", False, LIGHT_GRAY),
    ("  背驰信息 · 共振分析 · 模型切换", False, LIGHT_GRAY),
    ("", False, LIGHT_GRAY),
    ("路由：/stock/{code}", True, CYAN),
], 5.7, 1.7, 4.5, 2.8, font_size=11)

# 自选股页
add_rounded_rect(s8, 10.6, 1.15, 4.9, 3.6, CARD_BG)
add_rect(s8, 10.6, 1.15, 4.9, 0.45, ORANGE)
add_text(s8, "⭐  自选股页 /watchlist", 10.6, 1.2, 4.9, 0.4,
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_multiline_text(s8, [
    ("• 添加股票输入框", False, LIGHT_GRAY),
    ("• 自选股表格（代码/名称/现价/涨跌幅）", False, LIGHT_GRAY),
    ("• 点击行跳转个股分析页", False, LIGHT_GRAY),
    ("• 行末删除按钮", False, LIGHT_GRAY),
    ("• 持久化存储于 backend/watchlist.json", False, LIGHT_GRAY),
    ("", False, LIGHT_GRAY),
    ("路由：/watchlist", True, CYAN),
], 10.75, 1.7, 4.5, 2.8, font_size=11)

# K线可视化说明
add_rounded_rect(s8, 0.5, 4.9, 15, 2.5, CARD_BG)
add_text(s8, "📊  K 线可视化叠加层", 0.7, 5.0, 8, 0.4,
         font_size=14, bold=True, color=WHITE)
vis_items = [
    ("📈  MA均线", "MA5 · MA20 · MA60（独立开关）"),
    ("✏️  笔", "红涨绿跌实线（≥5根K线）"),
    ("📏  线段", "黄/橙粗线（≥3笔构成）"),
    ("🟣  中枢", "紫色半透明矩形+上下沿标注"),
    ("🎯  买卖点", "彩色圆点+文字标签"),
    ("🤖  AI信号", "入场线/止损线/止盈线"),
    ("📍  支撑阻力", "绿虚线（支撑）/ 红虚线（阻力）"),
]
for i, (title, desc) in enumerate(vis_items):
    col = i % 2 == 0
    x = 0.7 + (i // 2) * 5.1
    y = 5.42 + (i % 2) * 0.55
    add_text(s8, title, x, y, 2.2, 0.4, font_size=11, bold=True,
             color=GREEN if i % 2 == 0 else CYAN)
    add_text(s8, desc, x + 2.2, y, 2.6, 0.4, font_size=11, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════════
#  第 9 页：技术栈
# ════════════════════════════════════════════════════════════
s9 = add_slide()
fill_bg(s9, SLIDE_BG)
add_rect(s9, 0, 0, 16, 0.06, GREEN)

add_text(s9, "技术栈", 0.5, 0.25, 15, 0.7,
         font_size=32, bold=True, color=WHITE)
add_divider(s9, 0.5, 0.95, 2, GREEN, 0.03)

stacks = [
    ("后端", [
        ("Python", "3.10+", "主力语言", ORANGE),
        ("FastAPI", "0.115.x", "Web框架，异步API", GREEN),
        ("Uvicorn", "0.32.x", "ASGI服务器", TEAL),
        ("Pandas", "2.2.x", "K线数据处理", CYAN),
        ("AKShare", "1.14.x", "A股数据源", RED),
        ("httpx", "0.28.x", "HTTP客户端", RGBColor(0xA8, 0x55, 0xF7)),
        ("Pydantic", "2.10.x", "数据模型校验", RGBColor(0xF4, 0x72, 0x2E)),
    ]),
    ("前端", [
        ("Vue 3", "3.5.x", "渐进式JS框架", GREEN),
        ("TypeScript", "5.7.x", "类型安全", CYAN),
        ("Vite", "6.0.x", "构建工具", TEAL),
        ("Pinia", "2.3.x", "状态管理", ORANGE),
        ("Vue Router", "4.5.x", "前端路由", RGBColor(0xA8, 0x55, 0xF7)),
        ("ECharts", "5.5.x", "K线&副图图表库", RED),
        ("Axios", "1.7.x", "HTTP请求", RGBColor(0xF4, 0x72, 0x2E)),
    ]),
    ("AI模型（可选）", [
        ("DeepSeek", "API", "自然语言缠论分析", ORANGE),
        ("Gemini", "Google API", "自然语言缠论分析", GREEN),
    ]),
    ("辅助工具", [
        ("Git", "版本控制", "", CYAN),
        ("npm/pnpm", "包管理", "", TEAL),
        ("Playwright", "截图生成", "", RED),
    ]),
]

col_xs = [0.5, 5.5, 10.5, 13.5]
col_w = [4.8, 4.8, 2.8, 2.3]
for ci, (cat, items) in enumerate(stacks):
    x = col_xs[ci]; w = col_w[ci]
    add_rounded_rect(s9, x, 1.15, w, 6.7, CARD_BG)
    add_rect(s9, x, 1.15, w, 0.5, GREEN if ci < 2 else TEAL if ci == 2 else ORANGE)
    add_text(s9, cat, x + 0.15, 1.2, w - 0.2, 0.4,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        if len(item) == 4:
            name, ver, desc, col = item
        else:
            name, desc, _, col = item[0], item[1], "", CYAN
        y = 1.75 + j * 0.78
        add_rounded_rect(s9, x + 0.1, y, w - 0.2, 0.65,
                         RGBColor(0x0B, 0x10, 0x1A))
        add_text(s9, name, x + 0.2, y + 0.05, 2.0, 0.3,
                 font_size=11, bold=True, color=col)
        if len(item) >= 3:
            add_text(s9, ver, x + 0.2, y + 0.35, 2.0, 0.25,
                     font_size=9, color=MID_GRAY, font_name="Consolas")
        if desc:
            add_text(s9, desc, x + 0.2, y + 0.35, w - 0.4, 0.25,
                     font_size=9, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════════
#  第 10 页：常见问题 & 免责
# ════════════════════════════════════════════════════════════
s10 = add_slide()
fill_bg(s10, SLIDE_BG)
add_rect(s10, 0, 0, 16, 0.06, GREEN)

add_text(s10, "常见问题 & 免责声明", 0.5, 0.25, 15, 0.7,
         font_size=32, bold=True, color=WHITE)
add_divider(s10, 0.5, 0.95, 5, GREEN, 0.03)

qas = [
    ("Q: 后端报错 ModuleNotFoundError?", "确保已激活虚拟环境，执行 pip install -r requirements.txt"),
    ("Q: K线数据获取失败？", "AKShare依赖网络，部分接口有频率限制，稍后重试"),
    ("Q: SSL超时（Windows代理）？", "使用 python run_server.py 启动，自动处理网络代理兼容"),
    ("Q: 缠论分析结果为空？", "K线不足20根返回404，确认股票有足够交易历史"),
    ("Q: AI分析返回错误？", "检查.env中API Key配置及配额，可切换DeepSeek↔Gemini"),
    ("Q: 如何添加新指标？", "后端在kline_processor.py计算 → 前端新建XxxChart.vue → 引入StockView.vue"),
]
for i, (q, a) in enumerate(qas):
    row, ci = i // 2, i % 2
    x = 0.5 + ci * 7.7
    y = 1.15 + row * 1.4
    add_rounded_rect(s10, x, y, 7.4, 1.25, CARD_BG)
    add_text(s10, q, x + 0.15, y + 0.1, 7, 0.38, font_size=11, bold=True, color=ORANGE)
    add_text(s10, "A: " + a, x + 0.15, y + 0.52, 7, 0.6, font_size=11, color=LIGHT_GRAY)

# 免责
add_rounded_rect(s10, 0.5, 5.4, 15, 2.0, RGBColor(0x20, 0x0A, 0x0A))
add_rect(s10, 0.5, 5.4, 15, 0.06, RED)
add_text(s10, "⚠️  免责声明", 0.7, 5.55, 14.5, 0.45,
         font_size=16, bold=True, color=RED)
add_text(s10,
    "本系统仅供技术研究与学习使用，不构成任何投资建议。\n"
    "股票投资有风险，入市需谨慎。系统分析结果可能与实际走势存在偏差，"
    "投资者应自行承担决策风险。",
    0.7, 6.05, 14.5, 1.2, font_size=13, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════════
# 保存
# ════════════════════════════════════════════════════════════
output = "C:/code/stock-chanlun/docs/ChanStock_intro.pptx"
os.makedirs(os.path.dirname(output), exist_ok=True)
prs.save(output)
print(f"[OK] PPT saved: {output}")
print(f"     Total pages: {len(prs.slides)}")
